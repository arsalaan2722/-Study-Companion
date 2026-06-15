import streamlit as st
import os
import io
from dotenv import load_dotenv

load_dotenv()

from reportlab.lib.pagesizes import letter
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
from reportlab.lib.styles import getSampleStyleSheet
from pptx import Presentation
from pptx.util import Inches, Pt
from utils.pdf_reader import extract_text_from_pdf, extract_text_from_txt, extract_text_from_docx
from utils.summarizer import generate_summary
from utils.generator import generate_structured_notes

# Setup page config
st.set_page_config(page_title="AI Study Notes Generator", page_icon="🧠", layout="wide", initial_sidebar_state="expanded")

# Custom CSS for UI polishing
st.markdown("""
<style>
    /* Google Fonts */
    @import url('https://fonts.googleapis.com/css2?family=Inter:wght@400;500;600;700&display=swap');
    
    html, body, [class*="css"] {
        font-family: 'Inter', sans-serif;
    }
    
    /* Sidebar Logo */
    .sidebar-logo {
        color: #2563EB;
        font-weight: 700;
        font-size: 24px;
        margin-bottom: 2rem;
        display: flex;
        align-items: center;
        gap: 10px;
    }

    /* Output Section Styling */
    .report-badge {
        display: inline-block;
        background-color: #E0E7FF;
        color: #3730A3;
        padding: 4px 12px;
        border-radius: 16px;
        font-size: 12px;
        font-weight: 600;
        margin-bottom: 12px;
    }
    
    .report-title {
        font-size: 32px;
        font-weight: 700;
        color: #0F172A;
        margin-bottom: 24px;
        line-height: 1.2;
    }
    
    .summary-box {
        background-color: #EFF6FF;
        border-left: 4px solid #2563EB;
        padding: 24px;
        border-radius: 12px;
        margin-bottom: 24px;
        color: #334155;
        line-height: 1.6;
    }
    
    .summary-title {
        color: #1D4ED8;
        font-weight: 600;
        font-size: 18px;
        margin-bottom: 12px;
        display: flex;
        align-items: center;
        gap: 8px;
    }
    
    .content-card {
        background: white;
        padding: 32px;
        border-radius: 16px;
        margin-bottom: 24px;
        box-shadow: 0 4px 6px -1px rgba(0, 0, 0, 0.05);
        border: 1px solid #E2E8F0;
    }
    
    /* Make buttons full width */
    .stButton>button, .stDownloadButton>button {
        width: 100%;
        border-radius: 8px;
        font-weight: 600;
    }
</style>
""", unsafe_allow_html=True)

# Helper export functions
def create_pdf(summary, notes, key_points):
    buffer = io.BytesIO()
    doc = SimpleDocTemplate(buffer, pagesize=letter)
    styles = getSampleStyleSheet()
    Story = []
    Story.append(Paragraph("AI Generated Study Notes", styles['Title']))
    Story.append(Spacer(1, 12))
    Story.append(Paragraph("Summary", styles['Heading1']))
    Story.append(Paragraph(summary, styles['Normal']))
    Story.append(Spacer(1, 12))
    Story.append(Paragraph("Detailed Notes", styles['Heading1']))
    for line in notes.split('\n'):
        if line.strip():
            Story.append(Paragraph(line, styles['Normal']))
    Story.append(Spacer(1, 12))
    Story.append(Paragraph("Key Points", styles['Heading1']))
    for line in key_points.split('\n'):
        if line.strip():
            Story.append(Paragraph(line, styles['Normal']))
    doc.build(Story)
    buffer.seek(0)
    return buffer

def create_ppt(summary, notes, key_points):
    prs = Presentation()
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    slide.shapes.title.text = "AI Generated Study Notes"
    slide.placeholders[1].text = "Summary, Detailed Notes, and Key Points"
    
    bullet_slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(bullet_slide_layout)
    slide.shapes.title.text = "Summary"
    slide.placeholders[1].text_frame.text = summary
    
    for section in notes.split('\n\n'):
        if section.strip():
            slide = prs.slides.add_slide(bullet_slide_layout)
            slide.shapes.title.text = "Detailed Notes"
            tf = slide.placeholders[1].text_frame
            tf.word_wrap = True
            lines = section.split('\n')
            if lines[0].startswith('###') or lines[0].startswith('**'):
                slide.shapes.title.text = lines[0].replace('#', '').replace('*', '').strip()
                tf.text = '\n'.join(lines[1:])
            else:
                tf.text = section
                
    slide = prs.slides.add_slide(bullet_slide_layout)
    slide.shapes.title.text = "Key Points & Keywords"
    slide.placeholders[1].text_frame.text = key_points
    
    buffer = io.BytesIO()
    prs.save(buffer)
    buffer.seek(0)
    return buffer

if "generated_results" not in st.session_state:
    st.session_state.generated_results = None
if "current_page" not in st.session_state:
    st.session_state.current_page = "Home"
if "history" not in st.session_state:
    st.session_state.history = []
if "dark_mode" not in st.session_state:
    st.session_state.dark_mode = False

# Sidebar Content
with st.sidebar:
    st.markdown("<div class='sidebar-logo'>🧠 Study Companion</div>", unsafe_allow_html=True)
    st.markdown("<div style='color: #64748B; font-size: 12px; font-weight: 600; margin-bottom: 16px; letter-spacing: 1px;'>YOUR WORKSPACE</div>", unsafe_allow_html=True)
    
    if st.button("🏠 Home", use_container_width=True, type="primary" if st.session_state.current_page == "Home" else "secondary"):
        st.session_state.current_page = "Home"
        st.rerun()
    if st.button("✨ Features", use_container_width=True, type="primary" if st.session_state.current_page == "Features" else "secondary"):
        st.session_state.current_page = "Features"
        st.rerun()
    if st.button("🕒 History", use_container_width=True, type="primary" if st.session_state.current_page == "History" else "secondary"):
        st.session_state.current_page = "History"
        st.rerun()
    if st.button("ℹ️ About", use_container_width=True, type="primary" if st.session_state.current_page == "About" else "secondary"):
        st.session_state.current_page = "About"
        st.rerun()

    st.markdown("---")
    st.session_state.dark_mode = st.toggle("🌙 Dark Mode", value=st.session_state.dark_mode)

if st.session_state.dark_mode:
    st.markdown("""
    <style>
        .stApp { background-color: #000000 !important; }
        header[data-testid="stHeader"] { background-color: #000000 !important; }
        [data-testid="stSidebar"] { background-color: #0A0A0A !important; border-right: 1px solid #222222 !important; }
        .stMarkdown, .stText, h1, h2, h3, h4, h5, h6, p, label, span { color: #FFFFFF !important; }
        .content-card, .input-container { background-color: #0A0A0A !important; border-color: #222222 !important; box-shadow: none !important; }
        .summary-box { background-color: #000000 !important; border-left-color: #3B82F6 !important; }
        
        /* File Uploader */
        [data-testid="stFileUploadDropzone"] { background-color: #0A0A0A !important; border-color: #333333 !important; }
        [data-testid="stFileUploadDropzone"] * { color: #A3A3A3 !important; fill: #A3A3A3 !important; }
        .stFileUploader>div, .stFileUploader>div>div { background-color: #0A0A0A !important; border-color: #333333 !important; }
        .stFileUploader small, .stFileUploader [data-testid="stMarkdownContainer"] { color: #A3A3A3 !important; }
        
        /* Text Area */
        .stTextArea>div>div>textarea { background-color: #000000 !important; color: #FFFFFF !important; border: 1px solid #333333 !important; }
        .stTextArea>div>div>textarea::placeholder { color: #888888 !important; }
        
        /* Selectbox */
        .stSelectbox>div>div>div { background-color: #000000 !important; color: #FFFFFF !important; }
        
        /* Sidebar Buttons */
        button[kind="secondary"] { background-color: #111111 !important; color: #FFFFFF !important; border: 1px solid #333333 !important; }
        button[kind="secondary"]:hover { border-color: #3B82F6 !important; color: #3B82F6 !important; }
        button[kind="secondary"] p { color: #FFFFFF !important; }
        
        .report-title { color: #FFFFFF !important; }
        .sidebar-logo { color: #60A5FA !important; }
        .report-badge { background-color: #111827 !important; color: #60A5FA !important; border: 1px solid #1E3A8A !important; }
        .summary-title { color: #60A5FA !important; }
        .stDownloadButton>button { background-color: #111111 !important; color: #60A5FA !important; border-color: #60A5FA !important; }
        .stDownloadButton>button:hover { background-color: #111827 !important; }
        hr { border-color: #222222 !important; }
        [data-testid="stExpander"] { background-color: #0A0A0A !important; border-color: #222222 !important; }
    </style>
    """, unsafe_allow_html=True)

# Main layout
if st.session_state.current_page == "Home":
    col_left, col_spacing, col_right = st.columns([1, 0.1, 1.2])

    with col_left:
        st.markdown("<h2 style='font-size: 24px; color: #1E293B; margin-bottom: 20px;'>Source Material</h2>", unsafe_allow_html=True)
        
        st.markdown("<div class='input-container'>", unsafe_allow_html=True)
        uploaded_file = st.file_uploader("Drop PDF, DOCX, or TXT", type=["pdf", "docx", "txt"], help="Max file size: 10MB")
        
        text_input = ""
        if uploaded_file is not None:
            if uploaded_file.name.endswith(".pdf"):
                text_input = extract_text_from_pdf(uploaded_file)
            elif uploaded_file.name.endswith(".docx"):
                text_input = extract_text_from_docx(uploaded_file)
            elif uploaded_file.name.endswith(".txt"):
                text_input = extract_text_from_txt(uploaded_file)
                
        manual_text = st.text_area("", placeholder="Paste your topic or lecture notes here...", height=200, label_visibility="collapsed")
        if manual_text.strip():
            text_input = manual_text
        st.markdown("</div>", unsafe_allow_html=True)
        
        # Options Row
        opt_col1, opt_col2 = st.columns(2)
        with opt_col1:
            st.markdown("<div style='font-size: 12px; font-weight: 700; color: #2563EB; margin-bottom: 12px; letter-spacing: 0.5px;'>COMPONENTS</div>", unsafe_allow_html=True)
            gen_summary = st.checkbox("Summary", value=True)
            gen_detailed = st.checkbox("Detailed Notes", value=True)
            gen_keywords = st.checkbox("Key Points", value=True)
            
        with opt_col2:
            st.markdown("<div style='font-size: 12px; font-weight: 700; color: #2563EB; margin-bottom: 12px; letter-spacing: 0.5px;'>NOTES STYLE</div>", unsafe_allow_html=True)
            st.selectbox("", ["Academic Formal", "Simple & Student Friendly", "Bullet Point Heavy"], label_visibility="collapsed")
            
            st.write("") # Spacer
            if st.button("✨ Generate Notes", type="primary"):
                if not text_input.strip():
                    st.error("Please provide some text or upload a document.")
                elif not os.environ.get("GROQ_API_KEY"):
                    st.error("Please enter your Groq API key in the .env file.")
                else:
                    word_count = len(text_input.split())
                    if word_count > 5000:
                        text_input = " ".join(text_input.split()[:5000])
                        
                    with st.spinner("AI is analyzing your content..."):
                        try:
                            s = generate_summary(text_input) if gen_summary else ""
                            n = generate_structured_notes(text_input, task="notes") if gen_detailed else ""
                            k = generate_structured_notes(text_input, task="keywords") if gen_keywords else ""
                            
                            st.session_state.history.append({
                                "topic": text_input[:50],
                                "summary": s,
                                "notes": n,
                                "key_points": k
                            })
                            
                            st.session_state.generated_results = {
                                "topic": text_input[:50],
                                "summary": s,
                                "notes": n,
                                "key_points": k
                            }
                        except Exception as e:
                            st.error(f"Error: {str(e)}")

    with col_right:
        if st.session_state.generated_results:
            res = st.session_state.generated_results
            st.markdown("<div class='content-card'>", unsafe_allow_html=True)
            st.markdown("<span class='report-badge'>AI Generated Report</span>", unsafe_allow_html=True)
            
            words = res["topic"].split()
            title = " ".join(words[:5]).title() + ("..." if len(words) > 5 else "")
            st.markdown(f"<div class='report-title'>{title}</div>", unsafe_allow_html=True)
            
            if res["summary"]:
                st.markdown(f"<div class='summary-box'><div class='summary-title'>📄 Summary</div>{res['summary']}</div>", unsafe_allow_html=True)
                
            if res["notes"]:
                st.markdown("### 📝 Detailed Notes")
                st.markdown(res["notes"])
                st.markdown("---")
                
            if res["key_points"]:
                st.markdown("### 🔑 Key Points & Keywords")
                st.markdown(res["key_points"])
                
            st.markdown("</div>", unsafe_allow_html=True)
            
            d_col1, d_col2 = st.columns(2)
            with d_col1:
                pdf_file = create_pdf(res["summary"], res["notes"], res["key_points"])
                st.download_button(label="📄 Download PDF", data=pdf_file, file_name="study_notes.pdf", mime="application/pdf", type="primary")
            with d_col2:
                ppt_file = create_ppt(res["summary"], res["notes"], res["key_points"])
                st.download_button(label="📊 Download PPT", data=ppt_file, file_name="study_notes.pptx", mime="application/vnd.openxmlformats-officedocument.presentationml.presentation", type="primary")
        else:
            st.markdown("<div style='display: flex; height: 100%; align-items: center; justify-content: center; color: #94A3B8; flex-direction: column; padding-top: 100px;'><div style='font-size: 48px; margin-bottom: 16px;'>📝</div><h3>Your Notes Will Appear Here</h3><p>Upload a document or paste text to generate AI notes.</p></div>", unsafe_allow_html=True)

elif st.session_state.current_page == "Features":
    st.markdown("<h2 style='font-size: 32px; color: #1E293B; margin-bottom: 24px;'>✨ Features</h2>", unsafe_allow_html=True)
    st.markdown("<div class='content-card'><h3>📚 Smart Summarization</h3><p>Condense long lectures, articles, or books into clear, bite-sized summaries in seconds.</p><hr><h3>📝 Structured Note Generation</h3><p>Automatically turn messy text into highly organized, beautifully formatted notes with headings and bullet points.</p><hr><h3>🔑 Keyword Extraction</h3><p>Identify the most important terms and concepts automatically for rapid reviewing.</p><hr><h3>📄 Multi-format Support</h3><p>Drag and drop PDFs, Word Documents (.docx), or plain text files (.txt) directly into the app.</p><hr><h3>📊 One-Click Export</h3><p>Download your AI-generated study materials directly as PDF or PowerPoint presentations.</p></div>", unsafe_allow_html=True)

elif st.session_state.current_page == "History":
    st.markdown("<h2 style='font-size: 32px; color: #1E293B; margin-bottom: 24px;'>🕒 Your History</h2>", unsafe_allow_html=True)
    if len(st.session_state.history) == 0:
        st.info("You haven't generated any notes yet! Go to the Home tab to get started.")
    else:
        for i, item in enumerate(reversed(st.session_state.history)):
            with st.expander(f"Session {len(st.session_state.history) - i}: {item['topic'][:60]}..."):
                if item['summary']:
                    st.markdown("### Summary")
                    st.write(item['summary'])
                if item['notes']:
                    st.markdown("### Detailed Notes")
                    st.markdown(item['notes'])
                if item['key_points']:
                    st.markdown("### Key Points")
                    st.markdown(item['key_points'])

elif st.session_state.current_page == "About":
    st.markdown("<h2 style='font-size: 32px; color: #1E293B; margin-bottom: 24px;'>ℹ️ About</h2>", unsafe_allow_html=True)
    st.markdown("<div class='content-card'><p style='font-size: 16px; line-height: 1.6;'><b>Study Companion</b> is a powerful educational tool designed to help students, researchers, and professionals digest information faster.</p><p style='font-size: 16px; line-height: 1.6;'>Built with <b>Streamlit</b> and powered by the blazing-fast <b>Groq LLaMA 3.3 70B</b> AI model, it instantly transforms raw text and documents into organized, actionable knowledge.</p><p style='font-size: 16px; line-height: 1.6;'><i>Created to simplify learning and maximize productivity.</i></p></div>", unsafe_allow_html=True)

# Footer
st.markdown("---")
st.markdown("""
<div style='display: flex; justify-content: space-between; color: #64748B; font-size: 14px; padding: 20px 0;'>
    <div><strong>Study AI</strong> | Powered by AI</div>
    <div style='display: flex; gap: 20px;'>
        <span>Privacy Policy</span>
        <span>Contact Support</span>
        <span>Terms of Service</span>
    </div>
    <div>© 2026 Intelligent Learning Systems Inc.</div>
</div>
""", unsafe_allow_html=True)
